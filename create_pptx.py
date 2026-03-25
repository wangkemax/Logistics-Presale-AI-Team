#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
TITLE_COLOR = RGBColor(0, 51, 102)  # Dark blue
ACCENT_COLOR = RGBColor(204, 0, 0)  # Porsche red
TEXT_COLOR = RGBColor(51, 51, 51)   # Dark gray

def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        sub_frame = sub_box.text_frame
        sub_frame.word_wrap = True
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = TEXT_COLOR
        sub_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    return slide

def add_table_slide(prs, title, table_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Table
    rows = len(table_data)
    cols = len(table_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(12.333), Inches(rows * 0.5)).table
    
    for i, row in enumerate(table_data):
        for j, cell in enumerate(row):
            cell_obj = table.cell(i, j)
            cell_obj.text = cell
            if i == 0:  # Header row
                cell_obj.text_frame.paragraphs[0].font.bold = True
                cell_obj.text_frame.paragraphs[0].font.size = Pt(14)
                cell_obj.fill.solid()
                cell_obj.fill.fore_color.rgb = RGBColor(0, 51, 102)
                cell_obj.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            else:
                cell_obj.text_frame.paragraphs[0].font.size = Pt(12)
                cell_obj.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    
    return slide

# ============ SLIDES ============

# Slide 1: Title
add_title_slide(prs, "保时捷上海仓库项目投标方案", "核心内容与卖点共创会议")

# Slide 2: 会议背景与定位
add_content_slide(prs, "会议背景与定位", [
    "项目背景",
    "- 保时捷上海仓库项目已进入投标关键阶段",
    "- 方案整体框架已具备",
    "- 需要进一步统一并强化核心卖点与差异化表达",
    "",
    "本次会议定位",
    "- 不是重新设计整体方案",
    "- 而是：提炼最打动客户的核心内容，形成统一、清晰、可述标的卖点结构"
])

# Slide 3: 会议目标
add_content_slide(prs, "会议目标（Success Criteria）", [
    "1. 明确 3-5 个核心卖点",
    "   客户一听就懂，一听就记得住",
    "",
    "2. 为每个卖点形成清晰逻辑",
    "   - 客户痛点",
    "   - 我们的方案",
    "   - 核心差异",
    "   - 客户价值",
    "",
    "3. 明确卖点在投标中的使用方式",
    "   - PPT 所在章节",
    "   - 述标中的优先级"
])

# Slide 4: 会议讨论原则
add_content_slide(prs, "会议讨论原则", [
    "基本原则",
    "- 所有观点必须从客户视角出发",
    "- 不讨论价格与商务条款细节",
    "- 不追求面面俱到，而是重点突出",
    "- 本次会议允许、也必须淘汰卖点",
    "",
    "一句话共识",
    "我们不是在证明我们很强，",
    "而是在证明我们最适合保时捷。"
])

# Slide 5: 会议议程总览
add_table_slide(prs, "会议议程总览", [
    ["时间", "模块", "目标"],
    ["0-20 min", "客户视角对齐", "明确保时捷真正关注什么"],
    ["20-80 min", "结构化头脑风暴", "形成候选卖点"],
    ["80-120 min", "卖点筛选与排序", "确定核心卖点"],
    ["120-150 min", "落地与分工", "明确后续动作"]
])

# Slide 6: 模块一
add_content_slide(prs, "模块一：客户视角对齐", [
    "核心讨论问题",
    "",
    "1. 如果我们是保时捷，",
    "   为什么要选择/更换一个仓储与物流合作方？",
    "",
    "2. 在这个项目中，",
    "   保时捷最担心哪些风险？",
    "",
    "3. 对保时捷来说，",
    "   哪些能力是基本项，而不是卖点？",
    "",
    "本模块输出：形成 3-4 条客户核心关注点"
])

# Slide 7: 模块二
add_content_slide(prs, "模块二：结构化头脑风暴", [
    "一个合格卖点必须回答 4 个问题",
    "",
    "1. 客户的真实痛点是什么？",
    "2. 我们具体是如何解决的？",
    "3. 与其他 3PL 相比，差异在哪里？",
    "4. 这件事为保时捷带来的业务价值是什么？",
    "",
    "讨论方向示例",
    "- 运营稳定性与质量体系",
    "- 高端品牌匹配与形象",
    "- 数字化/合规/可视化",
    "- 柔性扩展与风险管理"
])

# Slide 8: 卖点记录模板
add_table_slide(prs, "卖点记录模板", [
    ["项目", "内容"],
    ["客户痛点", ""],
    ["我们的方案", ""],
    ["核心差异", ""],
    ["客户价值", ""],
    ["备注", ""]
])

# Slide 9: 模块三
add_content_slide(prs, "模块三：卖点筛选与排序", [
    "现实约束说明",
    "- 最终投标 PPT 中：核心卖点 <= 5 个，主卖点 <= 3 个",
    "",
    "筛选问题",
    "- 如果述标时间只有 10 分钟，哪些点必须讲？",
    "- 哪些点有价值，但不是第一优先级？",
    "",
    "筛选方法",
    "- 投票/打分",
    "- 强制排序"
])

# Slide 10: 模块四
add_content_slide(prs, "模块四：卖点如何落地到投标中", [
    "每个保留卖点需明确",
    "",
    "- 对应 PPT 章节",
    "- 是否为述标重点",
    "- 是否需要案例/数据/图示支撑",
    "- 后续深化责任人",
    "",
    "示例",
    "- 卖点A：述标主线（3-4 分钟）",
    "- 卖点B：重点亮点（1-2 分钟）",
    "- 卖点C：风险与稳定性补充点"
])

# Slide 11: 会议输出物确认
add_content_slide(prs, "会议输出物确认", [
    "本次会议最终输出",
    "",
    "1. 《核心卖点共识表》",
    "2. 卖点优先级排序结果",
    "3. 卖点深化与责任分工",
    "",
    "会后动作",
    "- 24 小时内输出会议纪要",
    "- 48 小时内完成投标方案初稿更新"
])

# Slide 12: 结束语
add_title_slide(prs, "会议结束语", "今天我们不是为了把方案做得更复杂，\n而是为了让客户更快理解我们、\n更加信任我们、更低风险地选择我们。")

# Save
output_path = os.path.expanduser("~/Desktop/保时捷上海仓库项目投标方案_核心卖点会议.pptx")
prs.save(output_path)
print("PPT saved to: " + output_path)
