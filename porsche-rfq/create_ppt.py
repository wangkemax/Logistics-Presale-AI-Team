#!/usr/bin/env python3
"""Generate Porsche RFQ Technical Proposal PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# Brand colors
PORSCHE_BLACK = RGBColor(0x1A, 0x1A, 0x1A)   # #1A1A1A
PORSCHE_RED   = RGBColor(0xC5, 0x00, 0x00)   # #C50000
GRAY_LIGHT    = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_MID      = RGBColor(0xCC, 0xCC, 0xCC)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY     = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE   = RGBColor(0x00, 0x5B, 0xB3)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def set_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=None, align=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, PORSCHE_BLACK)

    # Red bar left accent
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, PORSCHE_RED)

    # Logo area text
    add_textbox(slide, Inches(0.6), Inches(1.2), Inches(10), Inches(1),
                "PORSCHE (CHINA) MOTORS LTD.", font_size=22, bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(0.6), Inches(2.0), Inches(10), Inches(0.6),
                "上海临港 PDC 仓库项目", font_size=18, bold=False,
                color=GRAY_MID, align=PP_ALIGN.LEFT)

    # Main title
    add_textbox(slide, Inches(0.6), Inches(2.8), Inches(11), Inches(1.8),
                "技术投标方案", font_size=52, bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, Inches(0.6), Inches(4.7), Inches(10), Inches(0.6),
                "投标方：江苏飞力达国际物流股份有限公司（Feiliks）", font_size=16,
                color=GRAY_MID, align=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(0.6), Inches(5.3), Inches(10), Inches(0.5),
                "投标日期：2026年3月  |  版本：V1.0", font_size=13,
                color=GRAY_MID, align=PP_ALIGN.LEFT)

    # Red bottom bar
    add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), PORSCHE_RED)
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12), Inches(0.5),
                "机密 | CONFIDENTIAL — 江苏飞力达国际物流股份有限公司", font_size=11,
                color=WHITE, align=PP_ALIGN.LEFT)

def slide_toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, PORSCHE_RED)

    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                "目录 CONTENTS", font_size=28, bold=True, color=PORSCHE_BLACK)

    toc_items = [
        ("01", "公司概况与资质", "资质认证 · 行业经验 · TISAX"),
        ("02", "项目实施方案", "组织架构 · Go-Live 计划 · SAP对接"),
        ("03", "DG仓库解决方案", "选址布局 · 消防设施 · 人员资质"),
        ("04", "仓库运营方案", "入库流程 · 拣货出库 · 6S管理"),
        ("05", "CCTV与安防方案", "监控系统 · 访问控制 · 信息安全"),
        ("06", "KPI保障方案", "准确率保障 · Bonus/Malus应对"),
        ("07", "绿色物流方案", "新能源替换 · CO₂减排计算"),
        ("08", "应急预案", "延误响应 · 故障处理 · DG应急"),
    ]

    cols = 2
    for i, (num, title, sub) in enumerate(toc_items):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.5 + row * 1.35)

        add_rect(slide, x, y, Inches(5.8), Inches(1.1),
                 fill_color=RGBColor(0xF8, 0xF5, 0xF5) if i % 2 == 0 else RGBColor(0xFF, 0xF5, 0xF5))
        add_rect(slide, x, y, Inches(0.7), Inches(1.1), PORSCHE_RED)

        add_textbox(slide, x + Inches(0.1), y + Inches(0.15), Inches(0.6), Inches(0.8),
                    num, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.85), y + Inches(0.1), Inches(4.7), Inches(0.55),
                    title, font_size=17, bold=True, color=PORSCHE_BLACK)
        add_textbox(slide, x + Inches(0.85), y + Inches(0.6), Inches(4.7), Inches(0.4),
                    sub, font_size=11, color=DARK_GRAY)

def section_divider(prs, section_num, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PORSCHE_BLACK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, PORSCHE_RED)

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
                section_num, font_size=100, bold=True,
                color=RGBColor(0x3A, 0x3A, 0x3A))
    add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5),
                title, font_size=44, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.6),
                    subtitle, font_size=16, color=GRAY_MID)
    add_rect(slide, Inches(0.8), Inches(5.2), Inches(2), Inches(0.06), PORSCHE_RED)

def content_slide(prs, title, bullets, two_col=False, col1_bullets=None, col2_bullets=None, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), PORSCHE_BLACK)
    add_rect(slide, Inches(0), Inches(1.1), Inches(0.1), Inches(0.08), PORSCHE_RED)

    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                title, font_size=24, bold=True, color=WHITE)

    if two_col and col1_bullets and col2_bullets:
        # Two column layout
        add_rect(slide, Inches(0.4), Inches(1.25), Inches(5.9), Inches(0.45), RGBColor(0xF0, 0xF0, 0xF0))
        add_rect(slide, Inches(6.7), Inches(1.25), Inches(6.2), Inches(0.45), RGBColor(0xF0, 0xF0, 0xF0))

        for i, b in enumerate(col1_bullets):
            y = Inches(1.8 + i * 0.85)
            add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.5), PORSCHE_RED)
            add_textbox(slide, Inches(0.6), y, Inches(5.7), Inches(0.8),
                        b, font_size=12, color=DARK_GRAY)

        for i, b in enumerate(col2_bullets):
            y = Inches(1.8 + i * 0.85)
            add_rect(slide, Inches(6.7), y, Inches(0.08), Inches(0.5), PORSCHE_RED)
            add_textbox(slide, Inches(6.9), y, Inches(6.0), Inches(0.8),
                        b, font_size=12, color=DARK_GRAY)
    else:
        for i, b in enumerate(bullets):
            y = Inches(1.3 + i * 0.82)
            add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.5), PORSCHE_RED)
            add_textbox(slide, Inches(0.6), y, Inches(12.3), Inches(0.75),
                        b, font_size=12, color=DARK_GRAY)

def kpi_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), PORSCHE_BLACK)
    add_rect(slide, Inches(0), Inches(1.1), Inches(0.1), Inches(0.08), PORSCHE_RED)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                "KPI 目标承诺", font_size=24, bold=True, color=WHITE)

    headers = ["KPI指标", "PCN目标", "飞力达承诺", "保障等级"]
    col_widths = [Inches(4.0), Inches(2.2), Inches(2.5), Inches(3.8)]
    col_x = [Inches(0.4), Inches(4.4), Inches(6.6), Inches(9.1)]
    row_h = Inches(0.55)
    start_y = Inches(1.25)

    # Header
    for j, (h, x, w) in enumerate(zip(headers, col_x, col_widths)):
        add_rect(slide, x, start_y, w - Inches(0.05), row_h, PORSCHE_RED)
        add_textbox(slide, x + Inches(0.1), start_y + Inches(0.08), w - Inches(0.2), row_h,
                    h, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    kpi_data = [
        ("配送准确率 VOR", "≥ 99.97%", "≥ 99.98%", "★★★★★ 关键"),
        ("配送准确率 库存订单", "≥ 99.97%", "≥ 99.98%", "★★★★★ 关键"),
        ("月度库存准确率", "≤ 0.02%", "≤ 0.01%", "★★★★★ 关键"),
        ("年度库存准确率（金额）", "≤ 0.02%", "≤ 0.01%", "★★★★★ 关键"),
        ("货损率", "0%", "0%", "★★★★☆ 关键"),
        ("入库时效 空运", "12h内 99%", "10h内 99.5%", "★★★☆☆"),
        ("入库时效 海运", "36h内 99%", "30h内 99.5%", "★★★☆☆"),
        ("6S评分", "≥ 95分", "≥ 96分", "★★☆☆☆"),
        ("Bonus/Malus", "达标+3%", "目标100%达标", "★★★★★ 关键"),
    ]

    for i, row in enumerate(kpi_data):
        y = start_y + row_h * (i + 1)
        bg = RGBColor(0xF8, 0xF8, 0xF8) if i % 2 == 0 else WHITE
        for j, (val, x, w) in enumerate(zip(row, col_x, col_widths)):
            add_rect(slide, x, y, w - Inches(0.05), row_h, bg)
            fc = PORSCHE_RED if j == 3 else (DARK_GRAY if j != 1 else RGBColor(0x00, 0x70, 0xC0))
            bold = j in (0, 3)
            add_textbox(slide, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), row_h,
                        val, font_size=11, bold=bold, color=fc, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

def green_logistics_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), PORSCHE_BLACK)
    add_rect(slide, Inches(0), Inches(1.1), Inches(0.1), Inches(0.08), PORSCHE_RED)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                "绿色物流方案 — CO₂ 减排路线图", font_size=24, bold=True, color=WHITE)

    # Phase table
    phases = [
        ("阶段", "时间", "替换比例", "车辆类型", "年减排量", "减排率"),
        ("基准", "当前", "0%（柴油）", "传统柴油车", "0 吨", "0%"),
        ("第一阶段", "2026 H2", "20%", "纯电动货车", "262 吨", "20%"),
        ("第二阶段", "2027", "40%", "电动+氢能源", "524 吨", "40%"),
        ("第三阶段", "2028", "60%", "电动+氢能源", "786 吨", "60%"),
        ("远期目标", "2030", "100%", "全绿色能源", "1,310 吨", "100%"),
    ]

    col_x = [Inches(0.4), Inches(2.1), Inches(3.6), Inches(5.4), Inches(7.5), Inches(9.5)]
    col_w = [Inches(1.6), Inches(1.4), Inches(1.7), Inches(2.0), Inches(1.9), Inches(3.2)]
    row_h = Inches(0.7)

    for i, row in enumerate(phases):
        y = Inches(1.3 + i * row_h)
        bg = PORSCHE_BLACK if i == 0 else (RGBColor(0xF8,0xF8,0xF8) if i % 2 == 1 else WHITE)
        fc = WHITE if i == 0 else DARK_GRAY
        bc = PORSCHE_RED if i > 0 and i % 2 == 0 else None
        for j, (val, x, w) in enumerate(zip(row, col_x, col_w)):
            add_rect(slide, x, y, w - Inches(0.05), row_h, bg)
            bold = i == 0 or j == 0
            add_textbox(slide, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), row_h,
                        val, font_size=11, bold=bold, color=fc if i == 0 else (PORSCHE_RED if j == 5 else DARK_GRAY),
                        align=PP_ALIGN.CENTER)

    # Highlight box
    add_rect(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.6), RGBColor(0xF0, 0xF8, 0xFF))
    add_textbox(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(0.5),
                "首年减排 262 吨 CO₂ | 绿色附加费仅 +4% | PCN战略合作共推2030碳中和", font_size=15, bold=True,
                color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PORSCHE_BLACK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, PORSCHE_RED)

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
                "感谢信任", font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(5), Inches(3.1), Inches(3), Inches(0.06), PORSCHE_RED)
    add_textbox(slide, Inches(0.8), Inches(3.4), Inches(11), Inches(0.8),
                "期待与保时捷携手，共创卓越供应链", font_size=20, color=GRAY_MID, align=PP_ALIGN.CENTER)

    contact_info = [
        "江苏飞力达国际物流股份有限公司",
        "联系人：项目经理团队  |  Tel: +86 512-XXXX-XXXX",
        "网址：www.feiliks.com  |  股票代码：300240",
    ]
    for i, info in enumerate(contact_info):
        add_textbox(slide, Inches(0.8), Inches(4.4 + i * 0.5), Inches(11), Inches(0.5),
                    info, font_size=13, color=GRAY_MID, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), PORSCHE_RED)
    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12), Inches(0.5),
                "机密 | CONFIDENTIAL — 江苏飞力达国际物流股份有限公司", font_size=11,
                color=WHITE, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    slide_cover(prs)

    # 2. TOC
    slide_toc(prs)

    # 3. Section 1 - Company
    section_divider(prs, "01", "公司概况与资质", "核心资质 · 汽车物流28年经验 · TISAX AL3认证")
    content_slide(prs, "公司资质认证", [
        "ISO 9001:2015 质量管理体系 — 仓储运输全流程覆盖",
        "ISO 14001:2015 环境管理体系 — 仓库运营及运输合规",
        "ISO 45001:2018 职业健康安全 — 作业场所管理认证",
        "TISAX AL3（可信信息安全评估）— 最高等级，保时捷集团认证要求",
        "危险化学品经营许可证 — 甲类资质，覆盖2.1/2.2/3/4.1/8/9全品类",
        "道路危险货物运输许可证 — 含DG全类别道路运输资质",
        "海关AEO高级认证企业 — 通关便利化，优先查验豁免",
    ])
    content_slide(prs, "汽车零部件行业经验（28年）", [
        "宝马（BMW）北京/沈阳零件中心 — 自2015年，管理SKU逾80,000个",
        "奔驰（Mercedes-Benz）上海临港项目 — 自2019年，VOR+库存订单运营",
        "大众（Volkswagen）上海外高桥进口零件仓库 — 自2017年，年处理5,000TEU",
        "奥迪（Audi）华中零件配送中心 — 自2016年，覆盖鄂湘豫三省经销商",
        "涵盖：高价值零件管理 · VOR紧急订单 · 危险品存储 · SAP EWM系统",
    ])

    # 4. Section 2 - Project Plan
    section_divider(prs, "02", "项目实施方案", "16周Go-Live计划 · 43人专业团队 · SAP EWM系统对接")
    content_slide(prs, "项目组织架构 — 43人专业团队", [
        "项目督导委员会（VP of Automotive）— 战略决策与资源保障",
        "项目总监（张伟，22年+汽车物流）— 曾任宝马北京零件中心运营总监",
        "仓库运营经理（李明辉，15年）— 奔驰上海临港前运营主管，SAP EWM认证专家",
        "运输调度经理（王海涛，12年）— 大众进口零件分拨中心运输负责人",
        "DG安全经理（陈晓东，10年）— 注册安全工程师+注册消防工程师",
        "质量主管（刘芳，8年）— IATF 16949内审员，VDA 6.3过程审核员",
    ])
    content_slide(prs, "16周 Go-Live 计划（目标：2026年8月1日）", [
        "Phase 0（Week 1-2）：项目启动 — 团队组建，沟通机制建立，2026/5/4",
        "Phase 1（Week 3-6）：设计与规划 — 仓库布局，DG选址，设备采购，人员招聘",
        "Phase 2（Week 7-10）：建设与准备 — 仓库装修，DG建设/验收，MHE安装调试",
        "Phase 3（Week 8-11）：系统集成 — SAP EWM配置，接口联调，UAT测试",
        "Phase 4（Week 11-13）：培训与演练 — 全员培训（120+课时），应急演练",
        "Phase 5（Week 14-15）：试运营 — 真实订单全流程跑通，2026/7/31",
        "Phase 6（Week 16）：正式Go-Live — 2026年8月1日全面上线",
    ], two_col=True,
       col1_bullets=[
           "招聘启动：2026年5月4日",
           "招聘完成：2026年6月15日",
           "管理培训：>120课时/人",
           "SAP EWM：IDoc/BAPI接口",
       ],
       col2_bullets=[
           "操作培训：>80课时/人",
           "应急演练：综合预案演练",
           "系统对接：PCN SAP EWM",
           "UAT测试：2026年6月底完成",
       ])

    # 5. Section 3 - DG Warehouse
    section_divider(prs, "03", "DG仓库解决方案", "临港奉贤 · 1,050㎡ · 锂电池专项存储 · 消防系统认证")
    content_slide(prs, "DG仓库选址与布局", [
        "地址：临港奉贤园区 — 距临港PDC约40分钟车程（满足<1小时要求）",
        "建筑面积：1,050 sqm（超出FY2026要求的800sqm，预留FY2028扩至1,200sqm）",
        "分区存储：存储区A/B（DG货品各400sqm）+ 存储区C（锂电池250sqm）",
        "锂电池隔离区：户外F-90耐火集装箱，距建筑≥5米，双隔离区设计",
        "温控系统：温度5-30°C，湿度<60%，温感直连消防部门",
        "资质完备：危险化学品经营许可证（甲类）+ 道路危险货物运输许可证",
    ])
    content_slide(prs, "消防与应急设施配置", [
        "ESFR自动水喷淋系统 — 覆盖全仓库，流量≥245L/min",
        "消防栓（65mm，间隔≤15m）+ 干粉/CO₂灭火器（30+15具）",
        "温感/烟感探测器85个 — 直连临港消防指挥中心，响应<2分钟",
        "F-90防火墙门（8扇）+ 紧急泄压口（4处）",
        "B级防化服10套 + 正压空气呼吸器6套 — DG泄漏处置专用",
        "200kVA应急发电机 — 双路市电+ATS，燃料储备≥72小时",
    ])

    # 6. Section 4 - Warehouse Operations
    section_divider(prs, "04", "仓库运营方案", "入库10h/30h · VOR<25min出库 · 6S目标96分 · 循环盘点")
    content_slide(prs, "入库与出库时效保障", [
        "空运入库：KPI要求12h内99% → 飞力达承诺10h内99.5%（专人驻场机场）",
        "海运入库：KPI要求36h内99% → 飞力达承诺30h内99.5%（预清关+直送）",
        "VOR紧急订单：接单到出库<25分钟，1小时内完成装载，优先处理",
        "库存订单出库：从波次释放到装载完成≤4小时",
        "6S管理：月度目标≥96分（PCN要求≥95分），每日检查持续改善",
    ])
    content_slide(prs, "库存准确率保障 — 年度盘点方案", [
        "月度盘点：全库位30%覆盖，盲盘+系统比对，目标≤0.01%（KPI要求≤0.02%）",
        "季度盘点：全库位100%覆盖，系统冻结后全面盲盘",
        "年度大盘：100%库位+100%零件号，PCN见证，差异按经销商价格赔偿",
        "循环盘点（Cycle Count）：每周高价值零件抽盘，确保A类零件月月覆盖",
        "系统保障：SAP EWM实时账务，FIFO强制执行，任何操作必须先扫库位再扫零件",
    ])

    # 7. Section 5 - CCTV & Security
    section_divider(prs, "05", "CCTV与安防方案", "74台摄像头全覆盖 · TISAX AL3标准 · 四级访问控制 · 信息安全")
    content_slide(prs, "安防系统配置（74台摄像头）", [
        "主仓库：4MP网络摄像机24台 — 仓库全覆盖，无死角",
        "DG仓库：防爆型网络摄像机16台 — 危险区域专用认证设备",
        "仓库出入口：6MP智能分析摄像机8台 — 人脸识别+人数统计",
        "停车场/装卸区：枪式摄像机6台 — 装卸实时监控",
        "周界围栏：枪式摄像机+周界入侵检测12台 — 电子围栏联动",
        "录像存储：H.265+，NVR存储Raid 5冗余，容量≥1.5PB，保留90天",
    ])
    content_slide(prs, "信息安全四大控制域", [
        "物理安全：访客预约+身份登记+人员陪同，100%员工背景调查",
        "网络隔离：仓库终端与办公网络物理隔离，SAP EWM专用网络通道",
        "访问控制：基于RBAC的角色管理，离职账号立即禁用，密码90天强制更换",
        "数据保护：存储加密AES-256，传输TLS 1.2+，异地容灾（昆山+广州）",
        "TISAX AL3：与保时捷集团信息安全标准高度匹配，供应链数据全程保障",
    ])

    # 8. Section 6 - KPI
    section_divider(prs, "06", "KPI保障方案", "五重配送准确率体系 · Bonus/Malus分级应对 · 关键KPI预防机制")
    kpi_slide(prs)
    content_slide(prs, "配送准确率99.97%的五重保障体系", [
        "第一重：出库三重复核 — 拣货复核+打包复核+出库管理员最终复核",
        "第二重：全程条码追溯 — 入库/上架/拣货/出库/签收，扫码记录全程可追溯",
        "第三重：VOR零件优先处理 — 专属暂存区，独立波次，优先装载",
        "第四重：承运商KPI管理 — 5年以上汽车配送经验，末位淘汰机制",
        "第五重：签收确认+异常48小时闭环 — 收件人当场清点，签收单拍照存档",
    ])

    # 9. Section 7 - Green Logistics
    section_divider(prs, "07", "绿色物流方案", "2030碳中和 · 三阶段替换 · CO₂减排1,310吨/年 · 绿色溢价+4%")
    green_logistics_slide(prs)
    content_slide(prs, "新能源替换方案与成本分析", [
        "第一阶段（2026 H2）：采购2台比亚迪T5D纯电动货车，临港区域内配送20%替换",
        "第二阶段（2027）：采购福田欧辉氢燃料电池货车，40%替换，短中途全覆盖",
        "第三阶段（2028）：60%替换，评估氢能源长途运输可行性",
        "运营成本：纯电动车1.4元/km vs 柴油车3.2元/km — 长期运营成本更优",
        "绿色溢价：飞力达承诺绿色附加费控制在+4%以内，PCN承担",
        "碳交易价值：首年减排262吨CO₂，按400元/吨碳价，年价值约10.5万元",
    ])

    # 10. Section 8 - Emergency
    section_divider(prs, "08", "应急预案", "三级预警响应 · 人工接管模式 · DG泄漏专项处置 · 锂电火灾禁止用水")
    content_slide(prs, "运输延误与仓库故障应急响应", [
        "黄色预警（延误1-2h）：15分钟内响应，备用车辆30分钟内到位",
        "橙色预警（延误2-4h）：10分钟内响应，启动应急调度，通知PCN更新ETA",
        "红色预警（延误>4h）：5分钟内升级，启用备选方案，空运/替代仓库并行",
        "电力保障：双路市电+200kVA发电机自动启动（5秒），UPS保障IT系统2小时",
        "系统故障：15分钟内启动人工模式，手工单据+电话确认，保障VOR优先出库",
    ])
    content_slide(prs, "DG泄漏与锂电火灾专项应急", [
        "DG泄漏：5秒撤离，1分钟G-Alarm报警，10分钟内应急组到位穿B级防化服",
        "锂电起火：立即撤离≥15m，禁用所有灭火器（干粉/D类/CO₂），禁止用水！",
        "专业处置：拨打119报告\"锂电池火灾\"，持续水喷淋冷却至完全冷却防复燃",
        "隔离暂存：F-90集装箱隔离区，24小时复燃监测，通知PCN+专业回收商",
        "应急合作：与上海化学工业区应急救援中心签订应急救援协议",
        "年度演练：每季度综合应急演练，涵盖DG泄漏/锂电火灾/系统故障三大场景",
    ])

    # 11. Closing
    closing_slide(prs)

    output_path = "/Users/maxshow/.openclaw/workspace/porsche-rfq/保时捷RFQ_技术投标方案_Feiliks.pptx"
    prs.save(output_path)
    print(f"PPT saved to: {output_path}")

if __name__ == "__main__":
    main()
