#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Open template
template_path = '/Users/maxshow/Desktop/运输部分介绍.pptx'
prs = Presentation(template_path)

# Define content for each slide
slides_content = [
    {  # Slide 1 - Title
        "title": "保时捷上海仓库项目投标方案",
        "subtitle": "核心内容与卖点共创会议"
    },
    {  # Slide 2 - 会议背景与定位
        "title": "会议背景与定位",
        "bullets": [
            ("项目背景", [
                "保时捷上海仓库项目已进入投标关键阶段",
                "方案整体框架已具备",
                "需要进一步统一并强化核心卖点与差异化表达"
            ]),
            ("本次会议定位", [
                "不是重新设计整体方案",
                "提炼最打动客户的核心内容，形成统一、清晰、可述标的卖点结构"
            ])
        ]
    },
    {  # Slide 3 - 会议目标
        "title": "会议目标（Success Criteria）",
        "bullets": [
            ("1. 明确 3-5 个核心卖点", [
                "客户一听就懂，一听就记得住"
            ]),
            ("2. 为每个卖点形成清晰逻辑", [
                "客户痛点",
                "我们的方案",
                "核心差异",
                "客户价值"
            ]),
            ("3. 明确卖点在投标中的使用方式", [
                "PPT 所在章节",
                "述标中的优先级"
            ])
        ]
    },
    {  # Slide 4 - 会议讨论原则
        "title": "会议讨论原则",
        "bullets": [
            ("基本原则", [
                "所有观点必须从客户视角出发",
                "不讨论价格与商务条款细节",
                "不追求面面俱到，而是重点突出",
                "本次会议允许、也必须淘汰卖点"
            ]),
            ("一句话共识", [
                "我们不是在证明我们很强",
                "而是在证明我们最适合保时捷"
            ])
        ]
    },
    {  # Slide 5 - 会议议程 (need table)
        "title": "会议议程总览",
        "table": [
            ["时间", "模块", "目标"],
            ["0-20 min", "客户视角对齐", "明确保时捷真正关注什么"],
            ["20-80 min", "结构化头脑风暴", "形成候选卖点"],
            ["80-120 min", "卖点筛选与排序", "确定核心卖点"],
            ["120-150 min", "落地与分工", "明确后续动作"]
        ]
    },
    {  # Slide 6 - 模块一
        "title": "模块一：客户视角对齐",
        "bullets": [
            ("核心讨论问题", [
                "如果我们是保时捷，为什么要选择/更换一个仓储与物流合作方？",
                "在这个项目中，保时捷最担心哪些风险？",
                "对保时捷来说，哪些能力是基本项，而不是卖点？"
            ]),
            ("本模块输出", [
                "形成 3-4 条客户核心关注点"
            ])
        ]
    },
    {  # Slide 7 - 模块二
        "title": "模块二：结构化头脑风暴",
        "bullets": [
            ("一个合格卖点必须回答 4 个问题", [
                "客户的真实痛点是什么？",
                "我们具体是如何解决的？",
                "与其他 3PL 相比，差异在哪里？",
                "这件事为保时捷带来的业务价值是什么？"
            ]),
            ("讨论方向示例", [
                "运营稳定性与质量体系",
                "高端品牌匹配与形象",
                "数字化/合规/可视化",
                "柔性扩展与风险管理"
            ])
        ]
    },
    {  # Slide 8 - 卖点记录模板 (table)
        "title": "卖点记录模板",
        "table": [
            ["项目", "内容"],
            ["客户痛点", ""],
            ["我们的方案", ""],
            ["核心差异", ""],
            ["客户价值", ""],
            ["备注", ""]
        ]
    },
    {  # Slide 9 - 模块三
        "title": "模块三：卖点筛选与排序",
        "bullets": [
            ("现实约束说明", [
                "最终投标 PPT 中：核心卖点 <= 5 个，主卖点 <= 3 个"
            ]),
            ("筛选问题", [
                "如果述标时间只有 10 分钟，哪些点必须讲？",
                "哪些点有价值，但不是第一优先级？"
            ]),
            ("筛选方法", [
                "投票/打分",
                "强制排序"
            ])
        ]
    }
]

# Clear existing slides and add new content
# Keep only the first slide to use as template, delete others
while len(prs.slides) > 1:
    prs.slides.pop()

# Now fill in content
# For simplicity, let's just use the existing 9 slides and modify them
# Actually let's start fresh with blank slides

# Save
output_path = os.path.expanduser("~/Desktop/保时捷上海仓库项目投标方案_核心卖点会议_v2.pptx")
prs.save(output_path)
print("Template copied to: " + output_path)
